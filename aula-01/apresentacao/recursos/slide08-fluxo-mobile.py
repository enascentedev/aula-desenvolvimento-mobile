# -*- coding: utf-8 -*-
"""
Substitui o slide 08 (PWA) por "Como funciona uma aplicacao Mobile?".

Motivo: o slide 07 mostra o fluxo da Web (Usuario -> Navegador -> Web App ->
Servidor). O 08 passa a mostrar o fluxo equivalente do Mobile, criando o
paralelo estrutural direto. Fica mais alinhado a ementa, que pede a diferenca
entre nativo e hibrido, e nao exige PWA.

Reaproveita o mesmo vocabulario visual do slide 07: caixas ligadas por setas
vermelhas, primeira etapa em amarelo.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

VERMELHO = "D71920"
AMARELO_BG = "FFF9E6"
AMARELO_BD = "F1D777"
TINTA = "12161C"
SECUNDARIO = "5A6472"
SUPERFICIE = "F6F7F9"
BORDA = "E3E7EC"

BASE = os.path.dirname(os.path.abspath(__file__))
ICONES = os.path.join(BASE, "icones-alpha")


def texto_do(sh):
    return sh.text_frame.text.strip() if sh.has_text_frame else ""


def limpar_corpo(slide, manter_topo=1.45):
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        topo = Emu(sh.top).inches
        alt = Emu(sh.height).inches
        larg = Emu(sh.width).inches
        if larg < 0.6 and alt > 6:      # barras laterais
            continue
        if topo > 6.9:                  # rodape
            continue
        if topo >= manter_topo:
            sh._element.getparent().remove(sh._element)


def trocar_texto(shape, novo):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    modelo = p0.runs[0] if p0.runs else None
    nome = modelo.font.name if modelo else "Aptos"
    tam = modelo.font.size if modelo else Pt(14)
    try:
        cor = modelo.font.color.rgb
    except Exception:
        cor = RGBColor.from_string(TINTA)
    neg = modelo.font.bold if modelo else False

    tf.clear()
    for i, linha in enumerate(novo.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = linha
        r.font.name = nome
        r.font.size = tam
        r.font.color.rgb = cor
        r.font.bold = neg


def titulo(slide, texto):
    for sh in slide.shapes:
        if not texto_do(sh):
            continue
        top = Emu(sh.top).inches
        if 0.7 < top < 1.5 and Emu(sh.height).inches > 0.4:
            trocar_texto(sh, texto)
            return


def caixa(slide, x, y, w, h, texto, tam=12, cor=TINTA, fonte="Aptos",
          negrito=False, alinhar=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, espaco=None):
    cx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = cx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, linha in enumerate(texto.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinhar
        if espaco:
            p.space_after = Pt(espaco)
        r = p.add_run()
        r.text = linha
        r.font.size = Pt(tam)
        r.font.name = fonte
        r.font.bold = negrito
        r.font.color.rgb = RGBColor.from_string(cor)
    return cx


def cartao(slide, x, y, w, h, fundo=SUPERFICIE, borda=BORDA, raio=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(fundo)
    sh.line.color.rgb = RGBColor.from_string(borda)
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.adjustments[0] = raio
    if sh.has_text_frame:
        sh.text_frame.clear()
    return sh


def seta(slide, x, y, w):
    """Seta vermelha igual a do slide 07."""
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(x), Inches(y), Inches(w), Inches(0.12))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor.from_string(VERMELHO)
    s.line.fill.background()
    s.shadow.inherit = False
    if s.has_text_frame:
        s.text_frame.clear()
    return s


def montar(slide):
    limpar_corpo(slide)
    titulo(slide, "Como funciona uma aplicação Mobile?")
    caixa(slide, 0.95, 1.50, 10.5, 0.32,
          "Do toque do usuário até os dados.", tam=13, cor=SECUNDARIO)

    # ── Fluxo: mesmo vocabulario visual do slide 07 ──────────────────
    etapas = [
        ("Usuário", "toca no ícone", True),
        ("App instalado", "no dispositivo", False),
        ("Sistema\nOperacional", "Android / iOS", False),
        ("API", "requisição", False),
        ("Servidor", "dados", False),
    ]

    x = 0.95
    w = 1.92
    gap = 0.42
    y = 2.05
    h = 1.30

    for i, (nome, sub, destaque) in enumerate(etapas):
        if destaque:
            cartao(slide, x, y, w, h, fundo=AMARELO_BG, borda=AMARELO_BD)
        else:
            cartao(slide, x, y, w, h)

        caixa(slide, x, y + 0.32, w, 0.48, nome, tam=13.5, cor=TINTA,
              negrito=True, alinhar=PP_ALIGN.CENTER)
        caixa(slide, x, y + (0.86 if "\n" in nome else 0.78), w, 0.26, sub,
              tam=9.5, cor=SECUNDARIO, alinhar=PP_ALIGN.CENTER,
              fonte="Cascadia Mono")

        if i < len(etapas) - 1:
            seta(slide, x + w + 0.07, y + h / 2 - 0.06, gap - 0.14)

        x += w + gap

    # ── Detalhamento das etapas ──────────────────────────────────────
    # Sem quebra de linha forcada: o word wrap da caixa resolve melhor
    # e evita as quebras estranhas que apareceram na primeira versao.
    detalhes = [
        ("APP INSTALADO",
         "O usuário abre o aplicativo direto no aparelho. Não há URL e não há navegador."),
        ("CÓDIGO DA APLICAÇÃO",
         "Nativo: Kotlin / Java no Android, Swift no iOS. Multiplataforma: React Native, Flutter."),
        ("API",
         "Quando precisa de informação externa, o app faz uma requisição — igual à Web."),
        ("SERVIDOR / DADOS",
         "Processa, consulta o banco de dados e devolve uma resposta ao aplicativo."),
    ]

    dx = 0.95
    dw = 2.78
    dgap = 0.20
    dy = 3.70
    dh = 1.62

    for i, (rot, txt) in enumerate(detalhes):
        cartao(slide, dx, dy, dw, dh, fundo="FFFFFF", borda=BORDA)
        # numero da etapa
        selo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx + 0.22),
                                      Inches(dy + 0.22), Inches(0.30), Inches(0.30))
        selo.fill.solid()
        selo.fill.fore_color.rgb = RGBColor.from_string(VERMELHO)
        selo.line.fill.background()
        selo.shadow.inherit = False
        p = selo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.name = "Cascadia Mono"
        r.font.color.rgb = RGBColor.from_string("FFFFFF")

        caixa(slide, dx + 0.62, dy + 0.26, dw - 0.80, 0.28, rot,
              tam=9, cor=VERMELHO, fonte="Cascadia Mono", negrito=True)
        caixa(slide, dx + 0.24, dy + 0.72, dw - 0.46, 0.82, txt,
              tam=11, cor=SECUNDARIO)

        dx += dw + dgap

    # Icones das tecnologias no RODAPE do card 2 (no cabecalho eles
    # cobriam o rotulo "CODIGO DA APLICACAO").
    x_card2 = 0.95 + (dw + dgap)
    for j, ic in enumerate(["kotlin.png", "react.png"]):
        caminho = os.path.join(ICONES, ic)
        if os.path.exists(caminho):
            slide.shapes.add_picture(
                caminho,
                Inches(x_card2 + dw - 0.86 + j * 0.32),
                Inches(dy + dh - 0.44), Inches(0.24), Inches(0.24),
            )

    # ── Frase de fechamento ──────────────────────────────────────────
    cartao(slide, 0.95, 5.62, 11.44, 0.92, fundo=AMARELO_BG, borda=AMARELO_BD)
    caixa(slide, 1.22, 5.78, 11.0, 0.28, "A DIFERENÇA QUE IMPORTA",
          tam=9, cor="8A6D00", fonte="Cascadia Mono", negrito=True)
    caixa(slide, 1.22, 6.06, 11.0, 0.42,
          "Não precisamos abrir um navegador. O software está instalado no dispositivo "
          "e pode se integrar diretamente aos recursos do sistema operacional.",
          tam=12, cor=TINTA)


def main(entrada, saida):
    prs = Presentation(entrada)
    slide = list(prs.slides)[7]     # slide 08
    montar(slide)
    prs.save(saida)
    print(f"slide 08 substituido | salvo: {saida}")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
