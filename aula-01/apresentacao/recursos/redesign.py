"""
Redesign LIGHT tech - versao generica.

Detecta os "celulares" em QUALQUER slide pela geometria (corpo + tela + speaker)
e os substitui por aparelhos modernos com o icone da tecnologia na tela.

Identidade da Escola Tecnica Mesquita preservada (D71920 / FFD322 / logo).
"""
import os
import sys
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_FILL

TINTA = "12161C"
SECUNDARIO = "5A6472"
SUPERFICIE = "F6F7F9"
BORDA = "E3E7EC"
MOLDURA = "1B1F27"
TELA = "20242C"

BASE = os.path.dirname(os.path.abspath(__file__))
ICONES = os.path.join(BASE, "icones-alpha")

# Que icone usar conforme o texto que estava dentro/perto do celular
POR_TEXTO = {
    "WEB": "javascript.png",
    "NATIVO": "kotlin.png",
    "HÍBRIDO": "react.png",
    "ANDROID": "android.png",
    "IOS": "swift.png",
    "SITE": "javascript.png",
    "APP": "android.png",
    "WEB?": "javascript.png",
    "MULTI": "react.png",
}


def pol(v):
    return Emu(v).inches


def e_corpo_de_celular(sh):
    """Corpo de aparelho: retangulo alto, proporcao entre 1.6 e 2.4."""
    if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return False
    w, h = pol(sh.width), pol(sh.height)
    if w < 0.45 or w > 1.6:
        return False
    if h <= 0:
        return False
    r = h / w
    return 1.55 < r < 2.5


def dentro(a, b, folga=0.02):
    """a esta contido em b?"""
    fa = Inches(folga)
    return (a.left >= b.left - fa and a.top >= b.top - fa
            and a.left + a.width <= b.left + b.width + fa
            and a.top + a.height <= b.top + b.height + fa)


def desenhar_celular(slide, left, top, width, height, icone=None, rotulo=None, sub=None):
    moldura = max(Inches(0.045), int(width * 0.05))

    corpo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    corpo.fill.solid()
    corpo.fill.fore_color.rgb = RGBColor.from_string(MOLDURA)
    corpo.line.fill.background()
    corpo.shadow.inherit = False
    corpo.adjustments[0] = 0.14
    if corpo.has_text_frame:
        corpo.text_frame.clear()

    tela = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + moldura, top + moldura,
        width - 2 * moldura, height - 2 * moldura,
    )
    tela.fill.solid()
    tela.fill.fore_color.rgb = RGBColor.from_string(TELA)
    tela.line.fill.background()
    tela.shadow.inherit = False
    tela.adjustments[0] = 0.11
    if tela.has_text_frame:
        tela.text_frame.clear()

    ilha_w = int(width * 0.30)
    ilha = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + (width - ilha_w) // 2,
        top + moldura + Inches(0.04),
        ilha_w, Inches(0.07),
    )
    ilha.fill.solid()
    ilha.fill.fore_color.rgb = RGBColor.from_string(MOLDURA)
    ilha.line.fill.background()
    ilha.shadow.inherit = False
    ilha.adjustments[0] = 0.5
    if ilha.has_text_frame:
        ilha.text_frame.clear()

    if icone and os.path.exists(icone):
        ic = min(int(width * 0.48), Inches(0.60))
        slide.shapes.add_picture(
            icone,
            left + (width - ic) // 2,
            top + int(height * 0.28),
            ic, ic,
        )

    if rotulo:
        cx = slide.shapes.add_textbox(
            left - Inches(0.10), top + int(height * 0.57),
            width + Inches(0.20), Inches(0.24),
        )
        tf = cx.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = rotulo
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.name = "Aptos"
        r.font.color.rgb = RGBColor.from_string("FFFFFF")

    if sub:
        cx = slide.shapes.add_textbox(
            left - Inches(0.10), top + int(height * 0.70),
            width + Inches(0.20), Inches(0.20),
        )
        tf = cx.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = sub
        r.font.size = Pt(7.5)
        r.font.name = "Cascadia Mono"
        r.font.color.rgb = RGBColor.from_string("9AA4B2")


def refazer_celulares(slide):
    candidatos = [sh for sh in slide.shapes if e_corpo_de_celular(sh)]
    if not candidatos:
        return 0

    # A tela tambem tem proporcao de celular. Fica so o shape mais externo:
    # do maior para o menor, descarta todo candidato contido em outro ja aceito.
    candidatos.sort(key=lambda s: s.width * s.height, reverse=True)
    corpos = []
    for c in candidatos:
        if not any(dentro(c, j, folga=0.05) for j in corpos):
            corpos.append(c)

    trocados = 0
    for corpo in corpos:
        # Tudo que esta dentro do corpo sai junto
        internos = [sh for sh in slide.shapes if sh is not corpo and dentro(sh, corpo)]

        rotulo, sub = None, None
        for sh in internos:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip()
                if rotulo is None:
                    rotulo = t
                elif sub is None:
                    sub = t

        icone = None
        if rotulo:
            icone = POR_TEXTO.get(rotulo.strip().upper())
        if icone:
            icone = os.path.join(ICONES, icone)

        left, top = corpo.left, corpo.top
        width, height = corpo.width, corpo.height

        for sh in internos + [corpo]:
            try:
                slide.shapes._spTree.remove(sh._element)
            except Exception:
                pass

        # Proporcao de celular atual (~1:2.05) SEM crescer a caixa original:
        # mantem a altura e estreita a largura, recentralizando no lugar de antes.
        nova_w = int(height / 2.05)
        if nova_w > width:          # nunca alargar alem do espaco que havia
            nova_w = width
            height = int(width * 2.05)
        novo_left = left + (width - nova_w) // 2

        desenhar_celular(slide, novo_left, top, nova_w, height,
                         icone=icone, rotulo=rotulo, sub=sub)
        trocados += 1

    return trocados


def refinar_tema(slide):
    mapa = {
        "E7E4DE": SUPERFICIE,
        "FBFBFA": SUPERFICIE,
        "E5E7EB": BORDA,
        "1F1F1F": TINTA,
        "666666": SECUNDARIO,
        "5D6673": SECUNDARIO,
    }

    def trata(sh):
        try:
            if sh.fill.type == MSO_FILL.SOLID:
                h = str(sh.fill.fore_color.rgb).upper()
                if h in mapa:
                    sh.fill.fore_color.rgb = RGBColor.from_string(mapa[h])
        except Exception:
            pass
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        h = str(r.font.color.rgb).upper()
                        if h in mapa:
                            r.font.color.rgb = RGBColor.from_string(mapa[h])
                    except Exception:
                        pass
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s2 in sh.shapes:
                trata(s2)

    for sh in slide.shapes:
        trata(sh)


def main(entrada, saida):
    prs = Presentation(entrada)
    total = 0
    for i, s in enumerate(prs.slides, start=1):
        refinar_tema(s)
        n = refazer_celulares(s)
        if n:
            print(f"slide {i:2d}: {n} aparelho(s) redesenhado(s)")
            total += n
    prs.save(saida)
    print(f"total: {total} aparelhos | salvo: {saida}")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
