# -*- coding: utf-8 -*-
"""
Preenche as NOTAS DO APRESENTADOR de cada slide com o roteiro.

Reaproveita a mesma fonte de dados da pagina do roteiro visual
(gerar_pagina.SLIDES), para que apresentacao e roteiro nunca divirjam.

As notas aparecem no Modo de Exibicao do Apresentador do PowerPoint,
ao lado do cronometro e da miniatura do proximo slide.
"""
import sys
from pptx import Presentation
from pptx.util import Pt

from gerar_pagina import SLIDES, ROTULOS


# Ordem em que os blocos aparecem na nota. Fala primeiro: e o que o
# professor precisa achar em meio segundo enquanto fala.
PRIORIDADE = [
    "fala", "meta", "pergunta", "chave", "criterios",
    "foco", "anotar", "ordem", "planoA", "planoB", "dica", "design",
    "transicao", "atencao", "evitar", "novo", "corrigido",
]


def ordenar(itens):
    def chave(par):
        tipo = par[0]
        return PRIORIDADE.index(tipo) if tipo in PRIORIDADE else 99
    return sorted(itens, key=chave)


def montar_nota(num, titulo, tempo, bloco, itens):
    linhas = []
    linhas.append(f"SLIDE {num:02d} · {titulo}")
    linhas.append(f"[{tempo}] · {bloco}")
    linhas.append("")

    for tipo, texto in ordenar(itens):
        rotulo = ROTULOS.get(tipo, ("Nota", ""))[0].upper()
        # quebra em paragrafo proprio, com o rotulo destacado
        linhas.append(f"{rotulo}:")
        linhas.append(texto)
        linhas.append("")

    return "\n".join(linhas).rstrip()


def aplicar(entrada, saida):
    prs = Presentation(entrada)
    slides = list(prs.slides)

    if len(slides) != len(SLIDES):
        print(f"AVISO: {len(slides)} slides no arquivo, {len(SLIDES)} no roteiro")

    escritos = 0
    for num, titulo, tempo, bloco, itens in SLIDES:
        if num > len(slides):
            print(f"AVISO: slide {num} nao existe no arquivo")
            continue

        slide = slides[num - 1]
        nota = montar_nota(num, titulo, tempo, bloco, itens)

        tf = slide.notes_slide.notes_text_frame
        tf.clear()

        partes = nota.split("\n")
        for i, linha in enumerate(partes):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = linha
            r.font.size = Pt(14)          # legivel na tela do apresentador
            # cabecalho e rotulos em negrito
            if i < 2 or linha.endswith(":"):
                r.font.bold = True

        escritos += 1

    prs.save(saida)
    print(f"notas escritas em {escritos} slides | salvo: {saida}")


if __name__ == "__main__":
    aplicar(sys.argv[1], sys.argv[2])
